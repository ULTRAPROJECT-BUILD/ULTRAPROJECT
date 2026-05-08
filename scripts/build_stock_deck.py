#!/usr/bin/env python3
"""
Build a professional PowerPoint slide deck for stock research picks.
Ticket: T-019 / Project: stock-research-risky-plays
"""

import os
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn

# ── Color constants ──────────────────────────────────────────────────
NAVY       = RGBColor(0x1B, 0x2A, 0x4A)
CHARCOAL   = RGBColor(0x33, 0x33, 0x33)
STEEL_BLUE = RGBColor(0x44, 0x72, 0xC4)
LIGHT_GRAY = RGBColor(0xF2, 0xF2, 0xF2)
WATERMARK  = RGBColor(0xAA, 0xAA, 0xAA)
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
BLACK      = RGBColor(0x00, 0x00, 0x00)

DATE_STR = "2026-03-17"
FONT_NAME = "Calibri"
SLIDE_WIDTH  = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)
MARGIN = Inches(0.75)
CONTENT_WIDTH = SLIDE_WIDTH - 2 * MARGIN

# ── Stock data (extracted from the analysis markdown) ────────────────
stocks = [
    {
        "rank": 1, "ticker": "CRM", "name": "Salesforce", "sector": "Technology",
        "price": 195.31, "score": 69, "score_max": 100,
        "forward_pe": 13.1, "rsi": 53.1,
        "macd_signal": "Bullish (histogram +1.67)",
        "analyst_target": 276.33, "upside_pct": 41.5,
        "eps_forward": 14.91, "eps_trailing": 7.68,
        "revenue": "$41.5B", "revenue_growth": "Moderate (high-margin SaaS)",
        "profit_margin": "18%",
        "thesis": [
            "Beaten down 34% from 52-week high of $296.05, creating deep-value opportunity in a $183B market cap SaaS leader.",
            "Forward P/E of 13.1x is remarkably cheap vs tech sector avg ~22x (40% discount).",
            "MACD histogram turned positive at 1.67, signaling potential bullish crossover; price stabilized above 20-SMA."
        ],
        "risks": [
            "Broader tech selloff continuing amid macro uncertainty.",
            "AI spending cycle could slow enterprise software demand.",
            "Price remains below 50-SMA ($210.07), confirming intermediate-term downtrend."
        ],
        "allocation": 140, "alloc_pct": 14, "shares": 0.72,
        "entry": "$190-196", "stop": "$174.00", "t1": "$210.00", "t2": "$228.00",
    },
    {
        "rank": 2, "ticker": "ORCL", "name": "Oracle", "sector": "Technology",
        "price": 154.69, "score": 69, "score_max": 100,
        "forward_pe": 19.4, "rsi": 57.3,
        "macd_signal": "Bullish (histogram +1.58)",
        "analyst_target": 249.02, "upside_pct": 61.0,
        "eps_forward": 7.97, "eps_trailing": 5.52,
        "revenue": "$64.1B", "revenue_growth": "Strong (cloud transition)",
        "profit_margin": "25.3%",
        "thesis": [
            "Plunged 55% from 52-week high of $345.72 -- most dramatic dislocation in the screening universe.",
            "Analyst target of $249.02 implies 61% upside, the highest in the entire screen.",
            "Cloud infrastructure growth narrative intact; forward P/E 19.4x is reasonable for a growth play."
        ],
        "risks": [
            "Extreme debt/equity ratio of 415x (most leveraged in the screen).",
            "Sharp decline from $345 suggests possible structural concern, not just market weakness.",
            "Still below 50-SMA ($164.88), intermediate trend bearish."
        ],
        "allocation": 130, "alloc_pct": 13, "shares": 0.84,
        "entry": "$150-156", "stop": "$138.00", "t1": "$165.00", "t2": "$180.00",
    },
    {
        "rank": 3, "ticker": "AMD", "name": "Advanced Micro Devices", "sector": "Technology",
        "price": 196.31, "score": 62, "score_max": 100,
        "forward_pe": 18.3, "rsi": 39.3,
        "macd_signal": "Bullish (histogram +0.35)",
        "analyst_target": 290.27, "upside_pct": 47.9,
        "eps_forward": 10.75, "eps_trailing": 2.61,
        "revenue": "$34.6B", "revenue_growth": "Strong (AI GPU demand)",
        "profit_margin": "12.5%",
        "thesis": [
            "Down 26.5% from 52-week high of $267.08; RSI at 39.3 near oversold, suggesting a bounce setup.",
            "Forward P/E of 18.3x represents massive EPS growth from trailing 75.2x (312% EPS growth expected).",
            "Analyst target of $290.27 implies 47.9% upside; MACD histogram just turned positive."
        ],
        "risks": [
            "Semiconductor cycle risk -- AI capex may plateau.",
            "Competition from NVIDIA and Intel in data center GPU/CPU.",
            "Below both 20-SMA ($200.36) and 50-SMA ($215.23) -- bearish intermediate trend."
        ],
        "allocation": 150, "alloc_pct": 15, "shares": 0.76,
        "entry": "$192-198", "stop": "$186.00", "t1": "$210.00", "t2": "$230.00",
    },
    {
        "rank": 4, "ticker": "SLB", "name": "SLB Limited (fka Schlumberger)", "sector": "Energy / Oilfield Services",
        "price": 46.13, "score": 62, "score_max": 100,
        "forward_pe": 13.9, "rsi": 25.3,
        "macd_signal": "Bearish (histogram -0.61)",
        "analyst_target": 55.40, "upside_pct": 20.1,
        "eps_forward": 3.33, "eps_trailing": 2.41,
        "revenue": "$35.7B", "revenue_growth": "Moderate",
        "profit_margin": "9.5%",
        "thesis": [
            "Most technically oversold stock in screen (RSI 25.3) -- classic mean-reversion opportunity.",
            "Fell 12% from February highs near $52; bounced off Bollinger lower band with volume.",
            "Forward P/E of 13.9x is cheap for the oilfield services sector leader with $35.7B revenue."
        ],
        "risks": [
            "Oil price decline would hurt oilfield services spending.",
            "OPEC+ production cuts could unwind, reducing drilling activity.",
            "Price below both 20-SMA and 50-SMA -- bearish intermediate trend."
        ],
        "allocation": 140, "alloc_pct": 14, "shares": 3.03,
        "entry": "$44.50-46.50", "stop": "$42.50", "t1": "$49.00", "t2": "$52.00",
    },
    {
        "rank": 5, "ticker": "TSM", "name": "Taiwan Semiconductor", "sector": "Technology",
        "price": 345.98, "score": 56, "score_max": 100,
        "forward_pe": 19.3, "rsi": 29.8,
        "macd_signal": "Bearish (histogram -3.58)",
        "analyst_target": 430.65, "upside_pct": 24.5,
        "eps_forward": "N/A", "eps_trailing": "N/A",
        "revenue": "$3.81T TWD", "revenue_growth": "Strong",
        "profit_margin": "45.1%",
        "thesis": [
            "World's most critical semiconductor manufacturer ($1.79T market cap); oversold RSI at 29.8.",
            "Extraordinary 45.1% profit margin and 35.1% ROE make it a fundamental powerhouse.",
            "Analyst target of $430.65 implies 24.5% upside; Strong Buy consensus from 18 analysts."
        ],
        "risks": [
            "Geopolitical risk (Taiwan Strait tensions).",
            "Semiconductor cycle peak concerns.",
            "MACD strongly negative (histogram -3.58), momentum still bearish."
        ],
        "allocation": 140, "alloc_pct": 14, "shares": 0.40,
        "entry": "$338-348", "stop": "$325.00", "t1": "$360.00", "t2": "$385.00",
    },
    {
        "rank": 6, "ticker": "BIIB", "name": "Biogen", "sector": "Biotech",
        "price": 185.29, "score": 54, "score_max": 100,
        "forward_pe": 11.4, "rsi": 43.8,
        "macd_signal": "Bearish (histogram -1.09)",
        "analyst_target": 205.67, "upside_pct": 11.0,
        "eps_forward": 16.23, "eps_trailing": 8.89,
        "revenue": "$9.9B", "revenue_growth": "Moderate (legacy portfolio)",
        "profit_margin": "13.1%",
        "thesis": [
            "Forward P/E of 11.4x is the cheapest stock in the entire screen -- deep value pharma play.",
            "EPS expected to grow from $8.89 to $16.23 (82.6% growth); 30 analysts with Buy consensus.",
            "Stock sits above its 50-SMA ($184.25), providing a technical floor for entry."
        ],
        "risks": [
            "Alzheimer's drug (Leqembi) adoption slower than hoped.",
            "Revenue growth limited in legacy portfolio.",
            "RSI neutral at 43.8 -- no strong momentum signal."
        ],
        "allocation": 150, "alloc_pct": 15, "shares": 0.81,
        "entry": "$183-187", "stop": "$178.00", "t1": "$195.00", "t2": "$205.00",
    },
    {
        "rank": 7, "ticker": "REGN", "name": "Regeneron Pharmaceuticals", "sector": "Biotech",
        "price": 759.05, "score": 51, "score_max": 100,
        "forward_pe": 14.4, "rsi": 43.1,
        "macd_signal": "Bearish (histogram -3.10)",
        "analyst_target": 872.85, "upside_pct": 15.0,
        "eps_forward": 52.81, "eps_trailing": 41.60,
        "revenue": "$14.3B", "revenue_growth": "Solid",
        "profit_margin": "31.4%",
        "thesis": [
            "High-quality biotech with $14.3B revenue, 31.4% profit margins, and forward P/E of 14.4x.",
            "Pristine balance sheet with debt/equity of just 9.5x -- lowest leverage in biotech screen.",
            "Analyst target of $872.85 implies 15% upside; Buy consensus from 27 analysts."
        ],
        "risks": [
            "Eylea (key drug) facing biosimilar competition.",
            "RSI neutral with negative MACD histogram -- no momentum tailwind.",
            "Below both 20-SMA and 50-SMA -- intermediate trend bearish."
        ],
        "allocation": 150, "alloc_pct": 15, "shares": 0.20,
        "entry": "$750-765", "stop": "$738.00", "t1": "$790.00", "t2": "$820.00",
    },
]


# ── Helper functions ─────────────────────────────────────────────────

def set_font(run, size_pt, bold=False, italic=False, color=CHARCOAL, name=FONT_NAME):
    """Apply font settings to a run."""
    run.font.size = Pt(size_pt)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    run.font.name = name


def add_watermark(slide, date_str=DATE_STR):
    """Add small watermark text to the bottom-right corner."""
    left = SLIDE_WIDTH - Inches(3.3)
    top  = SLIDE_HEIGHT - Inches(0.4)
    txBox = slide.shapes.add_textbox(left, top, Inches(3.0), Inches(0.3))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.RIGHT
    run = p.add_run()
    run.text = f"AI-Generated Research | {date_str}"
    set_font(run, 8, color=WATERMARK, italic=True)


def add_title_textbox(slide, text, top, left=MARGIN, width=None, size=24, bold=True, color=NAVY):
    """Add a styled title / heading textbox."""
    if width is None:
        width = CONTENT_WIDTH
    txBox = slide.shapes.add_textbox(left, top, width, Inches(0.5))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    run = p.add_run()
    run.text = text
    set_font(run, size, bold=bold, color=color)
    return txBox


def add_bullet_list(slide, items, top, left=MARGIN, width=None, size=14, color=CHARCOAL, bold=False):
    """Add bullet-point text."""
    if width is None:
        width = CONTENT_WIDTH
    txBox = slide.shapes.add_textbox(left, top, width, Inches(len(items) * 0.35 + 0.1))
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.space_after = Pt(4)
        p.level = 0
        run = p.add_run()
        run.text = f"\u2022  {item}"
        set_font(run, size, bold=bold, color=color)
    return txBox


def set_cell(cell, text, size=12, bold=False, color=CHARCOAL, bg=None, alignment=PP_ALIGN.LEFT):
    """Set cell text and formatting."""
    cell.text = ""
    p = cell.text_frame.paragraphs[0]
    p.alignment = alignment
    run = p.add_run()
    run.text = str(text)
    set_font(run, size, bold=bold, color=color)
    # Reduce cell margins
    cell.margin_left = Inches(0.08)
    cell.margin_right = Inches(0.08)
    cell.margin_top = Inches(0.04)
    cell.margin_bottom = Inches(0.04)
    if bg:
        cell_fill = cell.fill
        cell_fill.solid()
        cell_fill.fore_color.rgb = bg


def add_table(slide, headers, rows, top, left=MARGIN, width=None, col_widths=None):
    """Add a styled table to a slide."""
    if width is None:
        width = CONTENT_WIDTH
    num_rows = len(rows) + 1  # +1 for header
    num_cols = len(headers)
    table_shape = slide.shapes.add_table(num_rows, num_cols, left, top, width, Inches(num_rows * 0.35))
    table = table_shape.table

    # Set column widths if provided
    if col_widths:
        for i, w in enumerate(col_widths):
            table.columns[i].width = w

    # Header row
    for j, header_text in enumerate(headers):
        set_cell(table.cell(0, j), header_text, size=12, bold=True, color=WHITE, bg=NAVY)

    # Data rows
    for i, row_data in enumerate(rows):
        bg = LIGHT_GRAY if i % 2 == 0 else WHITE
        for j, val in enumerate(row_data):
            set_cell(table.cell(i + 1, j), str(val), size=12, color=CHARCOAL, bg=bg)

    return table_shape


# ── Build the presentation ──────────────────────────────────────────

prs = Presentation()
prs.slide_width  = SLIDE_WIDTH
prs.slide_height = SLIDE_HEIGHT

# Use blank layout
blank_layout = prs.slide_layouts[6]  # Blank


# ────────────────────────────────────────────────────────────────────
# SLIDE 1: Title Slide
# ────────────────────────────────────────────────────────────────────
slide1 = prs.slides.add_slide(blank_layout)

# Background accent bar at top
bar = slide1.shapes.add_shape(
    1, Inches(0), Inches(0), SLIDE_WIDTH, Inches(0.15)  # 1 = rectangle
)
bar.fill.solid()
bar.fill.fore_color.rgb = NAVY
bar.line.fill.background()

# Title
add_title_textbox(slide1, "High-Risk / High-Upside NYSE Stock Picks",
                  top=Inches(2.0), size=32, color=NAVY)

# Subtitle
add_title_textbox(slide1, "AI-Generated Research",
                  top=Inches(2.8), size=24, bold=False, color=STEEL_BLUE)

# Date
add_title_textbox(slide1, f"Date: {DATE_STR}",
                  top=Inches(3.6), size=18, bold=False, color=CHARCOAL)

# Prepared for
add_title_textbox(slide1, "Prepared for Platform Admin",
                  top=Inches(4.2), size=18, bold=False, color=CHARCOAL)

# Bottom accent bar
bar2 = slide1.shapes.add_shape(
    1, Inches(0), SLIDE_HEIGHT - Inches(0.15), SLIDE_WIDTH, Inches(0.15)
)
bar2.fill.solid()
bar2.fill.fore_color.rgb = NAVY
bar2.line.fill.background()


# ────────────────────────────────────────────────────────────────────
# SLIDE 2: Executive Summary
# ────────────────────────────────────────────────────────────────────
slide2 = prs.slides.add_slide(blank_layout)
add_title_textbox(slide2, "Executive Summary", top=Inches(0.4), size=28)

summary_bullets = [
    "Methodology: 100-point scoring framework combining Technical (40%), Fundamental (30%), and Risk/Reward (30%) analysis.",
    "Universe: 15 NYSE-listed candidates screened across Technology, Biotech, and Energy sectors.",
    "Market Conditions: Tech names heavily sold off (20-55% declines from highs), creating deep-value entries; energy refiners overbought near highs.",
    "Result: 7 stock picks selected, led by CRM and ORCL tied at 69/100 points.",
    "Key Theme 1: Beaten-down technology names (CRM, ORCL, AMD) offer 41-61% analyst upside with improving technicals.",
    "Key Theme 2: Deeply oversold energy play (SLB, RSI 25.3) offers a classic mean-reversion trade setup.",
    "Key Theme 3: Biotech value plays (BIIB at 11.4x forward P/E, REGN with pristine balance sheet) provide defensive diversification.",
    "Budget: $1,000 fully allocated across 7 positions (max 15% per stock); time horizon 1-4 weeks; risk profile aggressive."
]
add_bullet_list(slide2, summary_bullets, top=Inches(1.1), size=14)
add_watermark(slide2)


# ────────────────────────────────────────────────────────────────────
# SLIDES 3-9: Individual Stock Pick Slides
# ────────────────────────────────────────────────────────────────────
for s in stocks:
    slide = prs.slides.add_slide(blank_layout)

    # Score label for risk/reward
    score_val = s["score"]
    # Convert 100-scale score to 10-scale for display
    score_10 = round(score_val / 10, 1)
    if score_10 >= 7.0:
        score_label = "Favorable"
    elif score_10 >= 5.0:
        score_label = "Moderate"
    else:
        score_label = "Cautious"

    # Header
    header_text = f"{s['name']}  ({s['ticker']})  |  ${s['price']:.2f}  |  {s['sector']}"
    add_title_textbox(slide, header_text, top=Inches(0.3), size=24, color=NAVY)

    # Score badge
    score_box = slide.shapes.add_textbox(SLIDE_WIDTH - Inches(3.0), Inches(0.3), Inches(2.25), Inches(0.45))
    tf = score_box.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.RIGHT
    run = p.add_run()
    run.text = f"Score: {score_val}/100"
    set_font(run, 20, bold=True, color=STEEL_BLUE)

    # ── Left column: Thesis + Risks ──
    left_col_left = MARGIN
    left_col_width = Inches(5.8)

    add_title_textbox(slide, "Investment Thesis", top=Inches(0.95), left=left_col_left,
                      width=left_col_width, size=16, color=STEEL_BLUE)
    add_bullet_list(slide, s["thesis"], top=Inches(1.35), left=left_col_left,
                    width=left_col_width, size=13)

    risk_top = Inches(3.0)
    add_title_textbox(slide, "Risk Factors", top=risk_top, left=left_col_left,
                      width=left_col_width, size=16, color=STEEL_BLUE)
    add_bullet_list(slide, s["risks"], top=risk_top + Inches(0.4), left=left_col_left,
                    width=left_col_width, size=13)

    # Risk/Reward score text
    rr_top = Inches(4.7)
    rr_box = slide.shapes.add_textbox(left_col_left, rr_top, left_col_width, Inches(0.4))
    tf = rr_box.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    run = p.add_run()
    run.text = f"Risk/Reward Score: {score_10}/10 -- {score_label}"
    set_font(run, 15, bold=True, color=NAVY)

    # ── Right column: Key Metrics Table ──
    right_col_left = Inches(7.0)
    right_col_width = Inches(5.5)

    add_title_textbox(slide, "Key Metrics", top=Inches(0.95), left=right_col_left,
                      width=right_col_width, size=16, color=STEEL_BLUE)

    # Determine EPS display
    eps_fwd = s["eps_forward"]
    eps_display = f"${eps_fwd}" if isinstance(eps_fwd, (int, float)) else str(eps_fwd)

    metrics_data = [
        ("P/E Ratio (Forward)", f"{s['forward_pe']}x"),
        ("RSI (14-day)", f"{s['rsi']}"),
        ("MACD Signal", s["macd_signal"]),
        ("Analyst Target", f"${s['analyst_target']:.2f}"),
        ("Upside %", f"{s['upside_pct']}%"),
        ("EPS (Forward)", eps_display),
        ("Revenue", s["revenue"]),
        ("Revenue Growth", s["revenue_growth"]),
        ("Profit Margin", s["profit_margin"]),
    ]

    metrics_headers = ["Metric", "Value"]
    metrics_rows = [[m[0], m[1]] for m in metrics_data]
    add_table(slide, metrics_headers, metrics_rows,
              top=Inches(1.35), left=right_col_left, width=right_col_width,
              col_widths=[Inches(2.5), Inches(3.0)])

    # ── Position sizing row at bottom ──
    pos_top = Inches(5.3)
    pos_items = [
        f"Allocation: ${s['allocation']} ({s['alloc_pct']}% of portfolio)  |  ~{s['shares']} shares  |  Entry: {s['entry']}  |  Stop: {s['stop']}  |  T1: {s['t1']}  |  T2: {s['t2']}"
    ]
    add_bullet_list(slide, pos_items, top=pos_top, size=12, color=CHARCOAL)

    add_watermark(slide)


# ────────────────────────────────────────────────────────────────────
# SLIDE 10: Comparison Table
# ────────────────────────────────────────────────────────────────────
slide10 = prs.slides.add_slide(blank_layout)
add_title_textbox(slide10, "Stock Picks at a Glance", top=Inches(0.4), size=28)

comp_headers = ["Ticker", "Company", "Price", "Sector", "P/E (Fwd)", "RSI", "Analyst Target", "Upside %", "Score"]
# Sort by score descending
sorted_stocks = sorted(stocks, key=lambda x: x["score"], reverse=True)
comp_rows = []
for s in sorted_stocks:
    comp_rows.append([
        s["ticker"],
        s["name"],
        f"${s['price']:.2f}",
        s["sector"],
        f"{s['forward_pe']}x",
        f"{s['rsi']}",
        f"${s['analyst_target']:.2f}",
        f"{s['upside_pct']}%",
        f"{s['score']}/100"
    ])

comp_table_shape = add_table(slide10, comp_headers, comp_rows, top=Inches(1.2),
                              col_widths=[Inches(0.9), Inches(2.8), Inches(1.0), Inches(2.0),
                                          Inches(1.1), Inches(0.8), Inches(1.5), Inches(1.2), Inches(1.0)])

# Bold the top pick rows (first two tied at 69)
table = comp_table_shape.table
for j in range(len(comp_headers)):
    cell = table.cell(1, j)
    for paragraph in cell.text_frame.paragraphs:
        for run in paragraph.runs:
            run.font.bold = True
    cell2 = table.cell(2, j)
    for paragraph in cell2.text_frame.paragraphs:
        for run in paragraph.runs:
            run.font.bold = True

# Add note about top picks
note_box = slide10.shapes.add_textbox(MARGIN, Inches(4.5), CONTENT_WIDTH, Inches(0.4))
tf = note_box.text_frame
p = tf.paragraphs[0]
run = p.add_run()
run.text = "Bold rows indicate top-ranked picks (CRM and ORCL tied at 69/100)."
set_font(run, 12, italic=True, color=STEEL_BLUE)

add_watermark(slide10)


# ────────────────────────────────────────────────────────────────────
# SLIDE 11: Position Sizing / Portfolio Allocation
# ────────────────────────────────────────────────────────────────────
slide11 = prs.slides.add_slide(blank_layout)
add_title_textbox(slide11, "Portfolio Allocation ($1,000 Budget)", top=Inches(0.4), size=28)

alloc_headers = ["#", "Ticker", "Company", "Sector", "Allocation", "% of Portfolio", "Shares (approx)", "Entry Zone", "Stop-Loss"]
alloc_rows = []
for s in stocks:
    alloc_rows.append([
        str(s["rank"]),
        s["ticker"],
        s["name"],
        s["sector"],
        f"${s['allocation']}",
        f"{s['alloc_pct']}%",
        f"{s['shares']:.2f}",
        s["entry"],
        s["stop"]
    ])
# Total row
alloc_rows.append(["", "TOTAL", "", "", "$1,000", "100%", "", "", ""])

add_table(slide11, alloc_headers, alloc_rows, top=Inches(1.2),
          col_widths=[Inches(0.4), Inches(0.8), Inches(2.5), Inches(2.0),
                      Inches(1.1), Inches(1.2), Inches(1.4), Inches(1.5), Inches(1.4)])

# Allocation notes
alloc_notes = [
    "Cash Reserve: $0 (fully allocated across 7 positions).",
    "Max single position: 15% ($150) -- no concentration risk.",
    "Sector diversification: 3 Technology, 2 Biotech, 1 Energy, 1 Oilfield Services.",
    "All positions use fractional shares; verify broker support before execution.",
    "Time horizon: 1-4 weeks (aggressive short-term plays)."
]
add_bullet_list(slide11, alloc_notes, top=Inches(4.7), size=13)
add_watermark(slide11)


# ────────────────────────────────────────────────────────────────────
# SLIDE 12: Disclaimer
# ────────────────────────────────────────────────────────────────────
slide12 = prs.slides.add_slide(blank_layout)
add_title_textbox(slide12, "Important Disclaimer", top=Inches(0.5), size=28)

disclaimer_lines = [
    "This report was generated by an AI system. It is not financial advice.",
    "",
    "The information contained herein is for informational and educational purposes only and should not be "
    "construed as a recommendation to buy, sell, or hold any security. Past performance does not guarantee "
    "future results. Always consult a licensed financial advisor before making investment decisions.",
    "",
    "The AI system that generated this report has no financial interest in any securities mentioned.",
    "",
    "Stock prices can decline to zero; never invest more than you can afford to lose. This analysis is based "
    f"on data available as of {DATE_STR} and may become outdated rapidly. The $1,000 budget and position "
    "sizing are illustrative; adjust based on your risk tolerance and financial situation.",
    "",
    "Data sources may have delays. Verify all figures before acting.",
    "",
    f"Data Sources: Yahoo Finance (via stock-data MCP), real-time quotes and 3-month daily OHLCV data.",
    f"Report Date: {DATE_STR}",
    f"Analysis generated by stock-analysis agent on {DATE_STR}."
]

disc_box = slide12.shapes.add_textbox(MARGIN, Inches(1.3), CONTENT_WIDTH, Inches(5.5))
tf = disc_box.text_frame
tf.word_wrap = True
for i, line in enumerate(disclaimer_lines):
    if i == 0:
        p = tf.paragraphs[0]
    else:
        p = tf.add_paragraph()
    p.alignment = PP_ALIGN.LEFT
    p.space_after = Pt(2)
    run = p.add_run()
    run.text = line
    set_font(run, 11, italic=True, color=CHARCOAL)

add_watermark(slide12)


# ────────────────────────────────────────────────────────────────────
# Save
# ────────────────────────────────────────────────────────────────────
repo_root = Path(__file__).resolve().parents[1]
output_dir = str(repo_root / "vault" / "clients" / "example" / "deliverables")
output_filename = f"stock-research-risky-plays-{DATE_STR}.pptx"
output_path = os.path.join(output_dir, output_filename)

os.makedirs(output_dir, exist_ok=True)
prs.save(output_path)

file_size = os.path.getsize(output_path)
slide_count = len(prs.slides)
tickers = [s["ticker"] for s in stocks]

print(f"pptx_path: {output_path}")
print(f"slide_count: {slide_count}")
print(f"picks_included: {', '.join(tickers)}")
print(f"file_size: {file_size:,} bytes ({file_size / 1024:.1f} KB)")
print("Done.")
